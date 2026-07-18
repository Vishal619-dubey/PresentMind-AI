import fitz,sys,os,tempfile,re,json
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
pdf,out=sys.argv[1],sys.argv[2];mode=sys.argv[3] if len(sys.argv)>3 else 'smart'
doc=fitz.open(pdf);prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5);tmp=tempfile.mkdtemp()

def clean(t): return re.sub(r'\s+',' ',t or '').strip()
def add_text(slide,text,x,y,w,h,size=18,bold=False,color=(45,50,75)):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=box.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=text;p.font.size=Pt(size);p.font.bold=bold;p.font.color.rgb=RGBColor(*color);return box

for i,page in enumerate(doc):
    pix=page.get_pixmap(matrix=fitz.Matrix(1.4,1.4),alpha=False);img=os.path.join(tmp,f'page-{i+1}.png');pix.save(img)
    slide=prs.slides.add_slide(prs.slide_layouts[6]);bg=slide.background.fill;bg.solid();bg.fore_color.rgb=RGBColor(248,249,253)
    if mode=='exact':
        slide.shapes.add_picture(img,0,0,width=prs.slide_width,height=prs.slide_height);continue
    blocks=[]
    for b in page.get_text('blocks'):
        txt=clean(b[4])
        if txt and len(txt)>1: blocks.append((b[1],txt))
    blocks.sort(key=lambda x:x[0]);texts=[t for _,t in blocks]
    title='Page '+str(i+1)
    for t in texts[:5]:
        if 4<=len(t)<=110:
            title=t;break
    add_text(slide,title,.65,.45,8.0,.8,27,True,(39,39,103))
    accent=slide.shapes.add_shape(1,Inches(.65),Inches(1.35),Inches(.9),Inches(.08));accent.fill.solid();accent.fill.fore_color.rgb=RGBColor(91,66,230);accent.line.color.rgb=RGBColor(91,66,230)
    slide.shapes.add_picture(img,Inches(9.15),Inches(1.15),width=Inches(3.55),height=Inches(5.5))
    body=[t for t in texts if t!=title and len(t)>15]
    body_text='\n'.join(body)
    sentences=re.split(r'(?<=[.!?])\s+|\n+',body_text)
    bullets=[]
    for s in sentences:
        s=clean(s)
        if 18<=len(s)<=240 and s not in bullets: bullets.append(s)
        if len(bullets)>=8: break
    if not bullets: bullets=['Text could not be reliably extracted from this page. Use Exact Visual mode for full fidelity.']
    box=slide.shapes.add_textbox(Inches(.75),Inches(1.55),Inches(7.9),Inches(5.25));tf=box.text_frame;tf.word_wrap=True
    for j,b in enumerate(bullets):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph();p.text=b;p.level=0;p.font.size=Pt(16 if len(bullets)>5 else 18);p.font.color.rgb=RGBColor(55,61,88);p.space_after=Pt(10)
    add_text(slide,f'Source page {i+1}',11.2,6.75,1.3,.25,9,False,(120,125,145))
prs.save(out)
