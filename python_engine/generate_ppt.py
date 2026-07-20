import json,sys
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
data=json.loads(sys.argv[1]);out=sys.argv[2]
prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)
for i,item in enumerate(data["outline"]):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background.fill;bg.solid();bg.fore_color.rgb=RGBColor(8,19,48) if i==0 else RGBColor(248,249,255)
    title=slide.shapes.add_textbox(Inches(.8),Inches(.7),Inches(11.7),Inches(1.2))
    p=title.text_frame.paragraphs[0];p.text=item["title"];p.font.size=Pt(30 if i else 38);p.font.bold=True;p.font.color.rgb=RGBColor(255,255,255) if i==0 else RGBColor(25,31,63)
    body=slide.shapes.add_textbox(Inches(.9),Inches(2),Inches(7.4),Inches(4.5))
    tf=body.text_frame
    if i==0:
        q=tf.paragraphs[0];q.text=item.get("summary","");q.font.size=Pt(20);q.font.color.rgb=RGBColor(190,205,240)
        r=tf.add_paragraph();r.text="Presented by Vishal Dubey";r.font.size=Pt(16);r.font.color.rgb=RGBColor(120,180,255)
    else:
        for j,b in enumerate(item.get("bullets",[])):
            q=tf.paragraphs[0] if j==0 else tf.add_paragraph();q.text=b;q.level=0;q.font.size=Pt(21);q.font.color.rgb=RGBColor(55,63,90);q.space_after=Pt(12)
        s=tf.add_paragraph();s.text=item.get("summary","");s.font.size=Pt(15);s.font.color.rgb=RGBColor(100,107,130)
    shape=slide.shapes.add_shape(1,Inches(9),Inches(1.8),Inches(3.3),Inches(3.3));shape.fill.solid();shape.fill.fore_color.rgb=RGBColor(74,66,230);shape.line.color.rgb=RGBColor(74,66,230)
    num=slide.shapes.add_textbox(Inches(10.1),Inches(2.8),Inches(1.2),Inches(1))
    pp=num.text_frame.paragraphs[0];pp.text=f"{i+1:02d}";pp.alignment=PP_ALIGN.CENTER;pp.font.size=Pt(38);pp.font.bold=True;pp.font.color.rgb=RGBColor(255,255,255)
prs.save(out)
